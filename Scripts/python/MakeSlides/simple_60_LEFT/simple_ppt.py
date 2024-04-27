'''
simple_ppt.py:   creates PPT slides using user inputs

Usage:     simply_ppt.py [-options] 'Song Title' [num of lines] [refrain]

options:
    -i:                            Italicize font

Arguments:
    Font size: [int]             Font size, default is 44
    Font color: (r ,g b)         Font color, default is black (0, 0 , 0)
'''

import os
import sys
import getopt
import collections
import collections.abc

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

# constants
TOP_ALIGN = 2.025
LEFT_ALIGN = -2.78
SLIDE_WIDTH = 30.96
SLIDE_HEIGHT = 15


def exit(msg=None):
    """Exit the program and print usage message"""
    print(__doc__)
    if msg:
        print(msg)
    sys.exit(0)


def create_slide(text_list, is_italic, font_size, font_color):
    """"These feature creates the slides"""
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = Cm(LEFT_ALIGN)
    top = Cm(TOP_ALIGN)
    width = Cm(SLIDE_WIDTH)
    height = Cm(SLIDE_HEIGHT)

    for line in text_list:
        txBox = slide.shapes.add_textbox(left, top, width, height)

        # txBox.fill.solid()
        # txBox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()

        run.text = f'''{line}'''
        if line == text_list[-1]:
            pass
        else:
            slide = prs.slides.add_slide(blank_slide_layout)

        font = run.font
        font.italic = is_italic
        font.name = 'Arial'
        font.size = Pt(font_size)
        # print(font_color[0])
        font.color.rgb = RGBColor(int(font_color[0]), int(
            font_color[1]), int(font_color[2]))

    prs.save('test.pptx')
    os.startfile('test.pptx')


# ----------- main function -----------
if __name__ == "__main__":
    try:
        opts, args = getopt.getopt(sys.argv[1:], 'i')
        numArgs = len(args)

        # assign default values
        is_italic = False
        font_size = 44
        font_color = (0, 0, 0)

        for opt, val in opts:
            if opt == '-i':
                is_italic = True

        if numArgs == 1:
            font_size = args[0]
        elif numArgs == 2:
            font_size, font_color = args

            try:
                font_size = int(font_size)

                font_color_temp = font_color.replace("'", '')
                font_color_temp = font_color_temp.replace(" ", '').split(',')

            except ValueError:
                exit('invalid parameters')
        else:
            exit()

        os.system(r'bullet_list.txt')

        with open("bullet_list.txt", "r", encoding="utf-8") as song_verse_data:
            txt_list = song_verse_data.read().split("\n")

    except FileNotFoundError as File_Error:
        print(File_Error)
    else:
        create_slide(txt_list, is_italic, font_size, font_color_temp)
