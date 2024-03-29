import os
import shutil
import collections
import collections.abc
from parameters import *

# Presentation is a function not a class
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from check_PPT_name import CheckSongName

class MakeSlide():
    def __init__(self, song_dict: dict, song_title: str, save_slide:bool, roadmap: list[str]):
        self.lyrics_dict = song_dict
        self.song_title = song_title
        self.save_slide = save_slide
        self.roadmap = roadmap

        #making new presentation
        self.prs = Presentation()
        self.blank_slide_layout = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(self.blank_slide_layout)


    def make_song_slides(self) -> None:
        # making title slide
        self.make_title_slides()

        # making verse slides
        for road_elm in self.roadmap:
            try:
                tmp_list = self.lyrics_dict[road_elm]
            except:
                print("roadmap element not it dictionary")
            else:
                italics = False
                if "Chorus" in road_elm or "Ending" in road_elm:
                    italics = True

                for line in tmp_list:
                    self.slide = self.prs.slides.add_slide(
                        self.blank_slide_layout)
                    self.make_blank_slide(
                        line, is_italic=italics, has_header=True, header=road_elm)

        pptx_name = f"{self.song_title}.pptx"
        self.prs.save(pptx_name)

        if self.save_slide:
            save_location = SONGS_DIRECTORY + pptx_name

            print(f"Saving {pptx_name} in {save_location}")
            
            if not os.path.exists(SONGS_DIRECTORY):
                os.mkdir(SONGS_DIRECTORY)
            shutil.copy2(pptx_name, save_location)

        os.startfile(f"{self.song_title}.pptx")

    def make_title_slides(self):
        self.make_blank_slide(self.song_title, is_title=True)
        check1 = CheckSongName(self.song_title)

        # if theres a match, exit code
        if check1.end_check:
            exit(0)

    def make_blank_slide(self, input_string, is_title=False, is_italic=False, has_header=False, header=None):
        txBox = self.slide.shapes.add_textbox(LEFT, TOP, WIDTH, HEIGHT)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()

        if is_title:
            run.text = f"""{input_string}"""
        else:
            run.text = f"""{input_string}"""

        font = run.font
        if is_title:
            font.bold = True
        if is_italic:
            font.italic = True
        font.name = "Calibri (Body)"
        # font.name = "Arial"
        font.size = Pt(70)

        if has_header:
            self.make_header(header)

    def make_header(self, header_str):
        # pass
        txBox = self.slide.shapes.add_textbox(
            LEFT, HEADER_TOP_ALIGN, WIDTH, HEADER_HEIGHT)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()

        run.text = header_str

        font = run.font
        font.name = "Calibri (Body)"
        font.italic = True
        font.size = Pt(24)
